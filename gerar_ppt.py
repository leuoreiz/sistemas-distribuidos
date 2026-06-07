# -*- coding: utf-8 -*-
"""Gera o PPT de 3 slides exigido no trabalho.

Conteúdo alinhado ao material do Prof. Carlos Kuretzki (linha Coulouris):
  - Premissas/características de sistemas distribuídos (A01 e A02)
  - Conceitos de nuvem: NIST, modelos de serviço/implantação (A02-C1, A04)
  - Vantagens da nuvem (A03)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL_CLARO = RGBColor(0x54, 0x68, 0xFF)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
FUNDO = RGBColor(0x0F, 0x12, 0x26)
CLARO = RGBColor(0xC8, 0xCD, 0xF0)
CINZA = RGBColor(0x9A, 0xA3, 0xD0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def fundo(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = FUNDO


def faixa(slide, titulo):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(12.1), Inches(1.0))
    tf = box.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = titulo
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = BRANCO
    ln = slide.shapes.add_textbox(Inches(0.65), Inches(1.08), Inches(5), Inches(0.1))
    rr = ln.text_frame.paragraphs[0].add_run(); rr.text = "—" * 20
    rr.font.color.rgb = AZUL_CLARO; rr.font.size = Pt(14); rr.font.bold = True


def bullets(slide, itens, top=1.6, left=0.7, width=12.1, size=16, gap=7):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.4))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for nivel, texto, negrito in itens:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = nivel; p.space_after = Pt(gap)
        marcador = "•  " if nivel == 0 else "    – "
        r = p.add_run(); r.text = marcador + texto
        r.font.size = Pt(size if nivel == 0 else size - 1)
        r.font.bold = negrito
        r.font.color.rgb = BRANCO if negrito else CLARO


def rodape(slide, texto):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.98), Inches(12.1), Inches(0.4))
    r = box.text_frame.paragraphs[0].add_run(); r.text = texto
    r.font.size = Pt(11); r.font.color.rgb = CINZA


# ---------- CAPA ----------
s = prs.slides.add_slide(BLANK); fundo(s)
t = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(2))
p = t.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Sistema Distribuído de Processamento de Tarefas"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BRANCO
sub = s.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(1.6))
for i, linha in enumerate([
    "Programação Distribuída / Computação em Nuvem — Prof. Carlos Kuretzki, PhD",
    "Aplicação de microsserviços em contêineres publicada em nuvem pública",
    "Integrantes: ____________________________________________",
]):
    p = sub.text_frame.paragraphs[0] if i == 0 else sub.text_frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = linha
    rr.font.size = Pt(16); rr.font.color.rgb = CINZA

# ---------- SLIDE 1 — Arquitetura e premissas ----------
s = prs.slides.add_slide(BLANK); fundo(s)
faixa(s, "1. Arquitetura e premissas de sistemas/programação distribuídos")
diag = s.shapes.add_textbox(Inches(0.7), Inches(1.42), Inches(12), Inches(0.6))
rr = diag.text_frame.paragraphs[0].add_run()
rr.text = "Navegador → Load Balancer (Nginx) → API (N réplicas) → Redis (fila + estado) → Workers (N)"
rr.font.size = Pt(14); rr.font.bold = True; rr.font.color.rgb = AZUL_CLARO
bullets(s, [
    (0, "Comunicação por troca de mensagens: API e workers só interagem via fila no Redis — a própria definição de SD do material.", True),
    (0, "Concorrência: vários workers processam tarefas em paralelo.", True),
    (0, "Escalabilidade: réplicas de API e de workers (replicação de serviços + caching no Redis).", True),
    (0, "Transparência (acesso, localização, replicação): o cliente usa um único endereço e não sabe qual nó respondeu.", True),
    (0, "Tolerância a falhas e confiabilidade: heartbeats com TTL; a tarefa permanece na fila se um nó cai (redundância).", True),
    (0, "Heterogeneidade: contêineres Docker executam em qualquer SO/ambiente.", True),
    (0, "Abertura (openness): API REST com interfaces publicadas (padrão aberto HTTP/JSON).", True),
    (0, "Compartilhamento de recursos + consistência: estado e contadores atômicos no Redis.", True),
    (0, "Acoplamento fraco e ausência de relógio global: nós stateless decidem com informação local.", True),
], top=2.0, size=15, gap=6)
rodape(s, "Stack: Python (FastAPI) • Redis • Nginx • Docker / docker-compose")

# ---------- SLIDE 2 — Escolha da nuvem e fundamentação ----------
s = prs.slides.add_slide(BLANK); fundo(s)
faixa(s, "2. Escolha da nuvem e fundamentação")
bullets(s, [
    (0, "Nuvem escolhida: Render — modelo de implantação NUVEM PÚBLICA.", True),
    (0, "Modelo de serviço: PaaS (Plataforma como Serviço) — fornece o ambiente de execução, sem gerenciar o servidor. Roda sobre a infraestrutura da AWS.", True),
    (0, "Atende às 5 características essenciais (NIST):", True),
    (1, "Auto-serviço sob demanda: deploy via Git, sem intervenção manual.", False),
    (1, "Acesso amplo via rede: aplicação acessível por HTTPS de qualquer lugar.", False),
    (1, "Agrupamento de recursos (pooling): infraestrutura compartilhada e gerenciada.", False),
    (1, "Rápida elasticidade: escala réplicas conforme a demanda.", False),
    (1, "Serviço de mensuração: uso medido (pay as you go).", False),
    (0, "Fundamentação prática: plano gratuito, deploy direto do Docker, HTTPS automático e infraestrutura como código (render.yaml).", True),
    (0, "Alternativas avaliadas: Amazon (AWS), Microsoft Azure e Google Cloud — equivalentes, porém exigem mais configuração e cartão de crédito.", False),
], top=1.55, size=15, gap=7)
rodape(s, "Deploy via Blueprint: New + → Blueprint → repositório → Apply  •  resultado: https://<app>.onrender.com")

# ---------- SLIDE 3 — Benefícios da solução ----------
s = prs.slides.add_slide(BLANK); fundo(s)
faixa(s, "3. Benefícios da solução implementada")
bullets(s, [
    (0, "Escalabilidade automática: adiciona nós/réplicas conforme a demanda, sem alterar o código.", True),
    (0, "Elasticidade: o ambiente aumenta ou diminui conforme a necessidade.", True),
    (0, "Alta disponibilidade e resiliência: redundância de nós + health checks + reinício automático; a fila preserva o trabalho se um nó falha.", True),
    (0, "Pay as you go: paga-se conforme o uso; no projeto, plano gratuito.", True),
    (0, "Portabilidade (Docker): a mesma imagem roda igual no ambiente local e na nuvem.", True),
    (0, "Desacoplamento: a fila absorve picos de carga; produtor e consumidor evoluem de forma independente.", True),
    (0, "Manutenção e evolução facilitadas: microsserviços pequenos, com deploy contínuo via Git.", True),
    (0, "Observabilidade: painel mostra fila, tarefas processadas e workers ativos em tempo real.", True),
], top=1.6, size=16, gap=11)
rodape(s, "Demonstração ao vivo: link na nuvem + escala/tolerância a falhas no docker-compose local")

prs.save("Trabalho_Computacao_em_Nuvem.pptx")
print("PPT gerado: Trabalho_Computacao_em_Nuvem.pptx")
